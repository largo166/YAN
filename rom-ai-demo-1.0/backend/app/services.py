import json
import re
from collections import Counter, defaultdict
from datetime import datetime
from pathlib import Path, PurePosixPath
from typing import Any, Optional

import httpx
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from sqlalchemy import delete, func, or_, select
from sqlalchemy.orm import Session

from app import models
from app.config import settings

SUPPORTED_PROJECT_EXTS = {".pdf", ".docx", ".xlsx", ".txt", ".md", ".pptx", ".png", ".jpg", ".jpeg"}
SUPPORTED_KNOWLEDGE_EXTS = {".md", ".txt", ".pdf", ".docx", ".xlsx", ".pptx", ".png", ".jpg", ".jpeg"}
MAX_BROWSER_UPLOAD_SIZE = 25 * 1024 * 1024
DEFAULT_VAULT_EXCLUDE_DIRS = {
    ".git",
    ".obsidian",
    ".venv",
    ".venv-funasr",
    "__pycache__",
    "node_modules",
    "dist",
    ".cache",
}
DEFAULT_VAULT_EXCLUDE_EXTS = {
    ".mp4",
    ".mov",
    ".avi",
    ".zip",
    ".rar",
    ".7z",
    ".psd",
    ".dwg",
    ".bak",
}
DEFAULT_VAULT_MAX_FILE_SIZE = 30 * 1024 * 1024

DIGITAL_EMPLOYEES = [
    ("AI项目经理", "项目负责人 / PM", ["任务拆解", "进度协调", "风险跟踪"], "PM"),
    ("AI总图助理", "总图负责人", ["强排", "流线", "指标复核"], "MP"),
    ("AI户型助理", "户型负责人", ["产品定位", "户型推演", "得房率"], "UN"),
    ("AI立面助理", "立面负责人", ["风格研究", "材料策略", "比例控制"], "FC"),
    ("AI景观助理", "景观接口", ["景观节点", "归家动线", "界面协同"], "LS"),
    ("AI室内助理", "室内接口", ["大堂", "会所", "精装氛围"], "IN"),
    ("AI汇报助理", "汇报助理", ["PPT目录", "汇报文案", "图像提示词"], "RP"),
    ("AI资料管理员", "资料管理员", ["文件整理", "知识库引用", "版本记录"], "KB"),
]

TEAM_MEMBERS = [
    ("项目经理", "项目负责人", ["业主沟通", "会议推进", "任务协调"], 45),
    ("主创负责人", "主创设计师", ["概念判断", "方案把控", "评审"], 55),
    ("总图负责人", "总图负责人", ["强排", "指标复核", "退界"], 50),
    ("立面负责人", "立面负责人", ["风格研究", "材料策略", "比例控制"], 35),
    ("汇报负责人", "汇报助理", ["PPT结构", "文本组织", "成果整合"], 40),
]


def ensure_digital_employees(db: Session) -> None:
    existing_names = set(db.scalars(select(models.DigitalEmployee.name)))
    target_names = {item[0] for item in DIGITAL_EMPLOYEES}
    if target_names.issubset(existing_names):
        return
    db.execute(delete(models.DigitalEmployee))
    for name, role, skills, avatar in DIGITAL_EMPLOYEES:
        db.add(
            models.DigitalEmployee(
                name=name,
                role=role,
                skills=json.dumps(skills, ensure_ascii=False),
                avatar=avatar,
                status="available",
                workload=20,
            )
        )
    db.commit()


def ensure_team_members(db: Session) -> None:
    existing_names = set(db.scalars(select(models.TeamMember.name)))
    if existing_names:
        return
    for name, role, skills, workload in TEAM_MEMBERS:
        db.add(
            models.TeamMember(
                name=name,
                role=role,
                skills=json.dumps(skills, ensure_ascii=False),
                status="available",
                workload=workload,
            )
        )
    db.commit()


def parse_document(path: Path) -> tuple[str, str]:
    ext = path.suffix.lower()
    try:
        if ext in {".png", ".jpg", ".jpeg"}:
            return "", "saved_no_ocr"
        if ext in {".txt", ".md"}:
            return path.read_text(encoding="utf-8", errors="ignore"), "parsed"
        if ext == ".pdf":
            reader = PdfReader(str(path))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
            return text, "parsed"
        if ext == ".docx":
            doc = Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
            return text, "parsed"
        if ext == ".xlsx":
            wb = load_workbook(str(path), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    values = [str(value) for value in row if value is not None]
                    if values:
                        rows.append(" | ".join(values))
            return "\n".join(rows), "parsed"
        if ext == ".pptx":
            try:
                from pptx import Presentation
            except Exception as exc:
                return f"解析失败：缺少 python-pptx（{exc}）", "failed"
            prs = Presentation(str(path))
            slides = []
            for idx, slide in enumerate(prs.slides, start=1):
                slides.append(f"# Slide {idx}")
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        slides.append(text)
            return "\n".join(slides), "parsed"
        return "", "unsupported"
    except Exception as exc:
        return f"解析失败：{exc}", "failed"


def project_upload_dir(project_id: str) -> Path:
    path = settings.upload_root_path / "projects" / project_id
    path.mkdir(parents=True, exist_ok=True)
    return path


def mock_analysis_payload(project: models.Project, files: list[models.ProjectFile]) -> dict[str, Any]:
    completeness = "低" if not files else "中"
    tasks = [
        {
            "task_name": "宋式展示区归家动线策略",
            "task_type": "策略分析",
            "priority": "高",
            "owner_role": "主创设计师",
            "estimated_days": 3,
            "dependencies": ["T001"],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "输出动线分析图、策略文本、汇报PPT结构。",
        },
        {
            "task_name": "社区礼序界面与门头尺度校准",
            "task_type": "立面研究",
            "priority": "高",
            "owner_role": "立面负责人",
            "estimated_days": 4,
            "dependencies": ["宋式展示区归家动线策略"],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "输出门头比例、材质建议和节点风险说明。",
        },
        {
            "task_name": "展示区景观节点与室内到访体验衔接",
            "task_type": "协同设计",
            "priority": "中",
            "owner_role": "景观接口",
            "estimated_days": 2,
            "dependencies": ["宋式展示区归家动线策略"],
            "risk_level": "低",
            "status": "todo",
            "output_requirement": "输出景观节点清单、室内接口条件和体验动线说明。",
        },
    ]
    timeline = [
        {
            "stage_name": "概念方案阶段",
            "start_day": 1,
            "end_day": 7,
            "milestone": "概念方案汇报",
            "dependencies": [],
            "risk_note": "甲方决策周期不确定。",
        },
        {
            "stage_name": "展示区深化阶段",
            "start_day": 8,
            "end_day": 14,
            "milestone": "展示区动线与界面确认",
            "dependencies": ["概念方案阶段"],
            "risk_note": "景观、室内、立面接口需要同步确认。",
        },
        {
            "stage_name": "汇报整合阶段",
            "start_day": 15,
            "end_day": 18,
            "milestone": "PPT与图像成果交付",
            "dependencies": ["展示区深化阶段"],
            "risk_note": "素材版本与甲方关注点可能发生变化。",
        },
    ]
    team_requirements = {
        "total_headcount": 5,
        "roles": [
            {"role": "项目负责人", "count": 1, "skills": ["项目管理", "甲方沟通"], "intensity": "全职"},
            {"role": "主创设计师", "count": 1, "skills": ["宋式风格", "展示区设计"], "intensity": "全职"},
            {"role": "立面负责人", "count": 1, "skills": ["比例控制", "材料策略"], "intensity": "阶段投入"},
            {"role": "景观接口", "count": 1, "skills": ["归家动线", "节点体验"], "intensity": "阶段投入"},
            {"role": "汇报助理", "count": 1, "skills": ["PPT结构", "图像提示词"], "intensity": "阶段投入"},
        ],
    }
    knowledge_refs = [
        {
            "source_file": "宋式社区方法论.md",
            "source_path": "Obsidian Vault/方法/宋式社区方法论.md",
            "chunk_id": "mock-song-community-001",
            "quote": "宋式展示区应强调归家动线的礼仪感，以门庭、院落、廊下空间形成连续体验。",
            "relevance_score": 0.92,
        },
        {
            "source_file": "展示区体验动线清单.md",
            "source_path": "Obsidian Vault/方法/展示区体验动线清单.md",
            "chunk_id": "mock-arrival-002",
            "quote": "首开展示区需要把车行到达、步行归家、接待转换和样板间参观整合为一条可讲述的路径。",
            "relevance_score": 0.86,
        },
    ]
    return {
        "mode": "mock",
        "project_basis": {
            "project_type": project.project_type or "待确认",
            "phase": project.phase or "待确认",
            "资料完整度": completeness,
            "缺失信息": ["规划条件细则", "成本目标", "甲方决策边界"] if not files else ["成本目标", "关键节点时间"],
            "综合风险等级": "中",
        },
        "design_difficulties": {
            "技术难点": ["需要尽快核对指标、总图边界、产品组合和消防/日照等基础约束。"],
            "协调难点": ["建筑、景观、室内和报规口径需要统一，避免汇报阶段反复返工。"],
            "规范难点": ["消防、日照、停车和无障碍需要在强排阶段提前校核。"],
            "甲方决策难点": ["需要明确产品档次、立面成本和展示区范围。"],
            "成本与落地难点": ["立面材料、景观节点和公区精装需要控制落地成本。"],
            "后续深化风险点": ["资料不完整会影响任务拆解、周期判断和团队配置。"],
        },
        "timeline_summary": {
            "总体推进周期": "约 28-42 天（Mock，需要人工校准）",
            "概念方案阶段": "5-7 天",
            "强排 / 总图阶段": "7-10 天",
            "户型与产品阶段": "5-8 天",
            "立面与风格阶段": "7-10 天",
            "景观 / 室内 / 精装协同阶段": "5-8 天",
            "报规或汇报成果阶段": "3-5 天",
            "关键路径": ["规划条件确认", "强排稳定", "立面方向锁定", "汇报成果整合"],
            "里程碑节点": ["概念方向会", "强排评审", "立面评审", "最终汇报"],
            "可能延误点": ["资料缺失", "甲方反复", "多专业接口不同步"],
        },
        "timeline": timeline,
        "staffing": {
            "项目负责人 / PM": "1人",
            "主创设计师": "1人",
            "总图负责人": "1人",
            "户型负责人": "1人",
            "立面负责人": "1人",
            "景观接口": "0.5人",
            "室内接口": "0.5人",
            "后期 / 报规接口": "0.5人",
            "AI辅助人员": "1人",
            "建议人数规模": "5-7人等效投入",
            "每类人员技能要求": ["强排经验", "产品定位", "立面落地", "汇报组织"],
            "各阶段人员投入强度": "前期 PM/总图高，汇报前主创/立面/汇报助理高。",
        },
        "team_requirements": team_requirements,
        "knowledge_refs": knowledge_refs,
        "tasks": tasks,
        "next_actions": [
            "补齐规划条件、红线、指标和甲方任务书。",
            "建立项目资料目录并标注版本。",
            "先做强排风险清单，再进入立面风格推演。",
            "明确展示区、首开区和汇报成果范围。",
            "把关键问题提交甲方形成一次决策会。",
        ],
    }


def analysis_to_markdown(payload: dict[str, Any]) -> str:
    lines = [f"# 项目分析报告（{payload.get('mode', 'mock')}）"]
    for section, value in payload.items():
        if section == "mode":
            continue
        lines.append(f"\n## {section}")
        if isinstance(value, dict):
            for key, item in value.items():
                rendered = json.dumps(item, ensure_ascii=False) if not isinstance(item, str) else item
                lines.append(f"- **{key}**：{rendered}")
        elif isinstance(value, list):
            for item in value:
                rendered = json.dumps(item, ensure_ascii=False) if isinstance(item, dict) else item
                lines.append(f"- {rendered}")
        else:
            lines.append(str(value))
    return "\n".join(lines)


def normalize_analysis_payload(payload: dict[str, Any], fallback: dict[str, Any]) -> dict[str, Any]:
    merged = {**fallback, **payload}
    for key in ("tasks", "timeline", "knowledge_refs"):
        if not isinstance(merged.get(key), list) or not merged[key]:
            merged[key] = fallback.get(key, [])
    if not isinstance(merged.get("team_requirements"), dict):
        merged["team_requirements"] = fallback.get("team_requirements", {"total_headcount": 0, "roles": []})
    return merged


async def call_deepseek_json(prompt: str) -> dict[str, Any]:
    if settings.mock_mode:
        return {}
    headers = {"Authorization": f"Bearer {settings.deepseek_api_key}", "Content-Type": "application/json"}
    body = {
        "model": settings.deepseek_model,
        "temperature": 0.2,
        "messages": [
            {
                "role": "system",
                "content": "你是建筑设计项目分析助手。只输出 JSON，不要输出 Markdown。字段必须稳定，缺失信息要明确标注。",
            },
            {"role": "user", "content": prompt},
        ],
    }
    async with httpx.AsyncClient(timeout=120) as client:
        response = await client.post(f"{settings.deepseek_base_url.rstrip('/')}/chat/completions", headers=headers, json=body)
        response.raise_for_status()
        text = response.json()["choices"][0]["message"]["content"]
    match = re.search(r"\{.*\}", text, re.S)
    if not match:
        return {}
    return json.loads(match.group(0))


async def call_deepseek_text(prompt: str, system_prompt: str) -> str:
    if settings.mock_mode:
        return ""
    headers = {"Authorization": f"Bearer {settings.deepseek_api_key}", "Content-Type": "application/json"}
    body = {
        "model": settings.deepseek_model,
        "temperature": 0.2,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": prompt},
        ],
    }
    async with httpx.AsyncClient(timeout=120) as client:
        response = await client.post(f"{settings.deepseek_base_url.rstrip('/')}/chat/completions", headers=headers, json=body)
        response.raise_for_status()
        return response.json()["choices"][0]["message"]["content"]


def save_analysis_result(db: Session, project: models.Project, payload: dict[str, Any], mode: str) -> models.ProjectReport:
    report = models.ProjectReport(
        project_id=project.id,
        report_type="project_analysis",
        content_json=json.dumps(payload, ensure_ascii=False),
        markdown=analysis_to_markdown(payload),
        model_name=settings.deepseek_model if mode == "deepseek" else "mock",
        mode=mode,
    )
    db.add(report)
    db.execute(delete(models.ProjectTask).where(models.ProjectTask.project_id == project.id))
    db.execute(delete(models.ProjectTimeline).where(models.ProjectTimeline.project_id == project.id))
    db.execute(delete(models.KnowledgeReference).where(models.KnowledgeReference.project_id == project.id))
    for item in payload.get("tasks", []):
        db.add(
            models.ProjectTask(
                project_id=project.id,
                task_name=item.get("task_name", ""),
                task_type=item.get("task_type", ""),
                priority=item.get("priority", "medium"),
                owner_role=item.get("owner_role", ""),
                estimated_days=int(item.get("estimated_days") or 1),
                dependencies=json.dumps(item.get("dependencies") or [], ensure_ascii=False),
                risk_level=item.get("risk_level", "medium"),
                status=item.get("status", "todo"),
                output_requirement=item.get("output_requirement", ""),
            )
        )
    for item in payload.get("timeline", []):
        db.add(
            models.ProjectTimeline(
                project_id=project.id,
                stage_name=item.get("stage_name", ""),
                start_day=int(item.get("start_day") or 1),
                end_day=int(item.get("end_day") or item.get("start_day") or 1),
                milestone=item.get("milestone", ""),
                dependencies=json.dumps(item.get("dependencies") or [], ensure_ascii=False),
                risk_note=item.get("risk_note", ""),
            )
        )
    for item in payload.get("knowledge_refs", []):
        db.add(
            models.KnowledgeReference(
                project_id=project.id,
                source_file=item.get("source_file", ""),
                source_path=item.get("source_path", ""),
                chunk_id=item.get("chunk_id", ""),
                quote=item.get("quote", ""),
                relevance_score=float(item.get("relevance_score") or 0),
            )
        )
    db.commit()
    db.refresh(report)
    return report


def ensure_project_sidecars(db: Session, project: models.Project) -> None:
    if project.tasks and project.timelines:
        return
    payload = mock_analysis_payload(project, project.files)
    save_analysis_result(db, project, payload, "mock")
    db.refresh(project)


def team_requirements_from_payload(payload: dict[str, Any]) -> dict[str, Any]:
    value = payload.get("team_requirements")
    if isinstance(value, dict):
        return value
    return {"total_headcount": 0, "roles": []}


def extract_md_metadata(text: str) -> tuple[str, list[str], list[str]]:
    heading = ""
    for line in text.splitlines():
        if line.startswith("#"):
            heading = line.lstrip("#").strip()
            break
    tags = sorted(set(re.findall(r"(?<!\w)#([\w\u4e00-\u9fff/-]+)", text)))
    links = sorted(set(re.findall(r"\[\[([^\]]+)\]\]", text)))
    return heading, tags, links


def chunk_text(text: str, path: str, tags: list[str], links: list[str]) -> list[dict[str, Any]]:
    if not text.strip():
        return []
    chunks = []
    current_heading = ""
    buffer = []
    for line in text.splitlines():
        if line.startswith("#") and buffer:
            content = "\n".join(buffer).strip()
            if content:
                chunks.append({"heading": current_heading, "content": content, "path": path, "tags": tags, "links": links})
            current_heading = line.lstrip("#").strip()
            buffer = [line]
        else:
            if line.startswith("#") and not current_heading:
                current_heading = line.lstrip("#").strip()
            buffer.append(line)
    content = "\n".join(buffer).strip()
    if content:
        chunks.append({"heading": current_heading, "content": content, "path": path, "tags": tags, "links": links})
    final = []
    for chunk in chunks:
        content = chunk["content"]
        if len(content) <= 1600:
            final.append(chunk)
            continue
        for i in range(0, len(content), 1400):
            final.append({**chunk, "content": content[i : i + 1600]})
    return final


def safe_browser_relative_path(filename: str) -> str:
    normalized = filename.replace("\\", "/").lstrip("/")
    parts = [part for part in PurePosixPath(normalized).parts if part not in {"", ".", ".."} and ":" not in part]
    return "/".join(parts) or "uploaded-file"


def index_knowledge_file(db: Session, source_path: Path, display_path: Optional[str] = None) -> dict[str, Any]:
    ext = source_path.suffix.lower()
    text, status = parse_document(source_path)
    heading, tags, links = extract_md_metadata(text) if ext == ".md" else ("", [], [])
    indexed_path = display_path or str(source_path)
    record = db.scalar(select(models.KnowledgeFile).where(models.KnowledgeFile.filepath == indexed_path))
    if not record:
        record = models.KnowledgeFile(filepath=indexed_path)
        db.add(record)
        db.flush()
    record.filename = Path(indexed_path).name
    record.filetype = ext.lstrip(".")
    record.filesize = source_path.stat().st_size
    record.title = heading or source_path.stem
    record.folder = str(Path(indexed_path).parent)
    db.execute(delete(models.KnowledgeChunk).where(models.KnowledgeChunk.file_id == record.id))
    db.execute(delete(models.KnowledgeTag).where(models.KnowledgeTag.file_id == record.id))
    db.execute(delete(models.KnowledgeLink).where(models.KnowledgeLink.file_id == record.id))
    for chunk in chunk_text(text, indexed_path, tags, links):
        db.add(
            models.KnowledgeChunk(
                file_id=record.id,
                heading=chunk["heading"],
                content=chunk["content"],
                path=chunk["path"],
                tags=json.dumps(chunk["tags"], ensure_ascii=False),
                links=json.dumps(chunk["links"], ensure_ascii=False),
            )
        )
    for tag, count in Counter(tags).items():
        db.add(models.KnowledgeTag(file_id=record.id, tag=tag, count=count))
    for link in links:
        db.add(models.KnowledgeLink(file_id=record.id, source_path=indexed_path, target=link))
    return {"filename": record.filename, "path": indexed_path, "filetype": record.filetype, "status": status}


def scan_knowledge_directory(db: Session, directory: Path, clear_existing: bool = False) -> dict[str, Any]:
    if clear_existing:
        db.execute(delete(models.KnowledgeLink))
        db.execute(delete(models.KnowledgeTag))
        db.execute(delete(models.KnowledgeChunk))
        db.execute(delete(models.KnowledgeFile))
        db.commit()
    stats = Counter()
    recent = []
    for path in directory.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in SUPPORTED_KNOWLEDGE_EXTS:
            continue
        ext = path.suffix.lower()
        stats["total_files"] += 1
        stats[ext.lstrip(".") or "other"] += 1
        recent.append(index_knowledge_file(db, path))
    db.commit()
    return {"indexed_files": stats["total_files"], "stats": dict(stats), "recent_files": recent[:20]}


def should_index_vault_path(path: Path, root: Path, include_sync_notes: bool = False) -> tuple[bool, str]:
    try:
        relative = path.relative_to(root)
    except ValueError:
        relative = path
    parts = set(relative.parts)
    if parts & DEFAULT_VAULT_EXCLUDE_DIRS:
        return False, "excluded_directory"
    if not include_sync_notes and "笔记同步助手" in parts:
        return False, "sync_notes_optional"
    if path.suffix.lower() not in SUPPORTED_KNOWLEDGE_EXTS:
        return False, "unsupported_type"
    if path.suffix.lower() in DEFAULT_VAULT_EXCLUDE_EXTS:
        return False, "excluded_type"
    try:
        if path.stat().st_size == 0:
            return False, "empty_file"
        if path.stat().st_size > DEFAULT_VAULT_MAX_FILE_SIZE:
            return False, "large_file"
    except OSError:
        return False, "stat_failed"
    return True, "ok"


def scan_vault_directory(
    db: Session,
    directory: Path,
    clear_existing: bool = False,
    include_sync_notes: bool = False,
) -> dict[str, Any]:
    if clear_existing:
        db.execute(delete(models.KnowledgeLink))
        db.execute(delete(models.KnowledgeTag))
        db.execute(delete(models.KnowledgeChunk))
        db.execute(delete(models.KnowledgeFile))
        db.commit()
    stats = Counter()
    skipped = Counter()
    recent = []
    priority_prefixes = ("wiki", "raw", "Obj-")
    for path in directory.rglob("*"):
        if not path.is_file():
            continue
        allowed, reason = should_index_vault_path(path, directory, include_sync_notes=include_sync_notes)
        if not allowed:
            skipped[reason] += 1
            continue
        try:
            rel = path.relative_to(directory)
            first = rel.parts[0] if rel.parts else ""
            priority = first == "wiki" or first == "raw" or first.startswith("Obj-")
        except ValueError:
            priority = False
        indexed = index_knowledge_file(db, path)
        indexed["priority_source"] = priority
        recent.append(indexed)
        stats["total_files"] += 1
        stats[path.suffix.lower().lstrip(".") or "other"] += 1
        if any(str(path).find(prefix) >= 0 for prefix in priority_prefixes):
            stats["priority_files"] += 1
    db.commit()
    return {
        "indexed_files": stats["total_files"],
        "stats": dict(stats),
        "skipped": dict(skipped),
        "recent_files": recent[:30],
        "filters": {
            "excluded_dirs": sorted(DEFAULT_VAULT_EXCLUDE_DIRS),
            "max_file_mb": DEFAULT_VAULT_MAX_FILE_SIZE // 1024 // 1024,
            "include_sync_notes": include_sync_notes,
        },
    }


def knowledge_stats(db: Session) -> dict[str, Any]:
    files = list(db.scalars(select(models.KnowledgeFile)))
    tags = list(db.scalars(select(models.KnowledgeTag)))
    links_count = db.scalar(select(func.count(models.KnowledgeLink.id))) or 0
    type_counts = Counter(file.filetype for file in files)
    top_tags = Counter()
    for tag in tags:
        top_tags[tag.tag] += tag.count
    return {
        "total_files": len(files),
        "markdown_files": type_counts.get("md", 0),
        "pdf_docx_xlsx_files": type_counts.get("pdf", 0) + type_counts.get("docx", 0) + type_counts.get("xlsx", 0),
        "image_files": type_counts.get("png", 0) + type_counts.get("jpg", 0) + type_counts.get("jpeg", 0),
        "filetype_distribution": dict(type_counts),
        "top_tags": [{"tag": tag, "count": count} for tag, count in top_tags.most_common(20)],
        "link_count": links_count,
    }


def knowledge_tree(db: Session) -> list[dict[str, Any]]:
    files = list(db.scalars(select(models.KnowledgeFile).order_by(models.KnowledgeFile.filepath)))
    roots: dict[str, Any] = {}
    for file in files:
        parts = Path(file.filepath).parts
        current = roots
        for part in parts[:-1]:
            current = current.setdefault(part, {})
        current.setdefault("_files", []).append({"name": file.filename, "path": file.filepath, "filetype": file.filetype})

    def pack(name: str, value: Any) -> dict[str, Any]:
        children = [pack(k, v) for k, v in value.items() if k != "_files"]
        files_nodes = [{"name": f["name"], "type": "file", "path": f["path"], "filetype": f["filetype"]} for f in value.get("_files", [])]
        return {"name": name, "type": "folder", "children": children + files_nodes}

    return [pack(k, v) for k, v in roots.items()]


def search_knowledge(db: Session, question: str, limit: int = 6) -> list[models.KnowledgeChunk]:
    terms = [term for term in re.split(r"\s+", question.strip()) if term]
    query = select(models.KnowledgeChunk)
    if terms:
        clauses = []
        for term in terms[:6]:
            clauses.append(models.KnowledgeChunk.content.ilike(f"%{term}%"))
            clauses.append(models.KnowledgeChunk.heading.ilike(f"%{term}%"))
        query = query.where(or_(*clauses))
    return list(db.scalars(query.limit(limit)))


def _markdown_list(items: list[str]) -> str:
    return "\n".join(f"- {item}" for item in items)


def _refs_from_chunks(chunks: list[models.KnowledgeChunk]) -> list[dict[str, Any]]:
    refs = []
    for chunk in chunks:
        refs.append(
            {
                "chunk_id": chunk.id,
                "source_file": Path(chunk.path).name,
                "source_path": chunk.path,
                "heading": chunk.heading,
                "quote": chunk.content[:420],
                "relevance_score": 0.82,
            }
        )
    return refs


def build_startup_analysis_payload(project: models.Project, chunks: list[models.KnowledgeChunk]) -> dict[str, Any]:
    refs = _refs_from_chunks(chunks)
    ref_names = [ref["source_file"] for ref in refs[:6]]
    description = project.description or "暂无项目描述，需由项目经理补充业主诉求、规模、阶段和交付目标。"
    technical_focus_cards = [
        {
            "title": "技术卡_日照计算",
            "dimension": "日照",
            "summary": "启动阶段应确认当地日照规则、分析口径、遮挡边界和是否已有日照复核成果。",
            "checkpoints": ["大寒日/冬至日口径", "被遮挡对象", "楼栋间距", "复核图纸版本"],
            "source_refs": ref_names,
            "manual_confirm": "需要人工确认当地规划部门采用的日照计算口径。",
        },
        {
            "title": "技术卡_退界要求",
            "dimension": "退界",
            "summary": "优先核对红线、道路、邻地、消防登高面及地下边界退距。",
            "checkpoints": ["用地红线", "道路退距", "邻地退距", "地下室边界", "消防登高面"],
            "source_refs": ref_names,
            "manual_confirm": "需要以规划条件或设计任务书为准。",
        },
        {
            "title": "技术卡_面积计算",
            "dimension": "面积",
            "summary": "明确计容/不计容、赠送面积、架空层、人防、地下室和配套用房计算方式。",
            "checkpoints": ["计容总建面", "可售面积", "地下室", "架空层", "配套用房", "人防面积"],
            "source_refs": ref_names,
            "manual_confirm": "需要和当地面积测绘/报规口径交叉确认。",
        },
        {
            "title": "技术卡_消防风险",
            "dimension": "消防",
            "summary": "强排前置检查消防车道、登高场地、间距、防火分区和地下车库组织。",
            "checkpoints": ["消防车道", "登高场地", "防火间距", "地下车库", "疏散组织"],
            "source_refs": ref_names,
            "manual_confirm": "需要专业负责人复核。",
        },
        {
            "title": "技术卡_规划条件",
            "dimension": "规划",
            "summary": "项目启动必须锁定容积率、限高、密度、绿地率、停车、配套和地块边界。",
            "checkpoints": ["容积率", "限高", "建筑密度", "绿地率", "停车配比", "配套要求"],
            "source_refs": ref_names,
            "manual_confirm": "缺规划条件时，不应进入最终任务拆解。",
        },
        {
            "title": "技术卡_报批风险",
            "dimension": "报批",
            "summary": "识别需要提前沟通的规划、消防、人防、面积、日照和专家会风险。",
            "checkpoints": ["规划沟通", "消防审查", "人防口径", "面积测算", "日照复核", "专家会"],
            "source_refs": ref_names,
            "manual_confirm": "需要项目经理形成待业主确认清单。",
        },
    ]
    task_breakdown = [
        {
            "task_name": "项目启动资料完整性检查",
            "task_type": "启动检查",
            "priority": "高",
            "owner_role": "项目经理",
            "estimated_days": 1,
            "dependencies": [],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "形成资料缺口表，确认规划条件、红线、任务书和历史参考资料。",
        },
        {
            "task_name": "历史项目技术复用提取",
            "task_type": "知识复用",
            "priority": "高",
            "owner_role": "AI资料管理员",
            "estimated_days": 1,
            "dependencies": ["项目启动资料完整性检查"],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "输出日照、退界、面积、消防、规划、报批六类技术卡。",
        },
        {
            "task_name": "启动会议程与业主确认清单",
            "task_type": "会议推进",
            "priority": "高",
            "owner_role": "项目经理",
            "estimated_days": 1,
            "dependencies": ["历史项目技术复用提取"],
            "risk_level": "低",
            "status": "todo",
            "output_requirement": "形成启动会议程、待确认问题和下一步责任人。",
        },
        {
            "task_name": "项目启动汇报PPT结构",
            "task_type": "汇报组织",
            "priority": "中",
            "owner_role": "AI汇报助理",
            "estimated_days": 2,
            "dependencies": ["启动会议程与业主确认清单"],
            "risk_level": "中",
            "status": "todo",
            "output_requirement": "输出面向业主的PPT目录和每页核心表达。",
        },
    ]
    meeting_agenda = [
        "确认项目背景、区位、规模、设计阶段和业主诉求。",
        "核对规划条件、红线、日照、退界、面积计算和消防/报批资料缺口。",
        "复用历史项目经验，明确哪些规则可参考、哪些必须重新确认。",
        "确认本轮交付物、责任人和时间节点。",
        "形成下次会议前的待办和业主决策事项。",
    ]
    ppt_outline = [
        {"page": 1, "title": "项目启动背景", "content": "项目基本信息、业主诉求、当前阶段。"},
        {"page": 2, "title": "资料完整度与缺口", "content": "已上传资料、缺失资料、需确认来源。"},
        {"page": 3, "title": "历史项目技术复用", "content": "日照、退界、面积、消防、规划、报批六类重点。"},
        {"page": 4, "title": "启动阶段任务拆解", "content": "任务、负责人、优先级、交付物。"},
        {"page": 5, "title": "风险与下一步", "content": "未决问题、业主决策事项、下一次会议。"},
    ]
    risk_list = [
        "规划条件、红线或任务书缺失会导致任务拆解偏差。",
        "日照、退界、面积计算口径必须以当地正式要求为准。",
        "历史项目只能作为复用参考，不能替代本项目审批依据。",
        "若会议结论不回写任务看板，项目经理后续追踪会断链。",
    ]
    open_questions = [
        "业主本轮最关心的是进度、产品定位、成本还是报批风险？",
        "是否已有正式规划条件、红线和设计任务书？",
        "本项目是否需要先做强排可行性或日照快速复核？",
        "启动会后哪些事项需要业主书面确认？",
    ]
    mindmap_json = {
        "title": "项目启动分析",
        "nodes": [
            {
                "id": "project",
                "label": project.name,
                "children": [
                    {
                        "id": "technical",
                        "label": "技术重点",
                        "children": [{"id": card["dimension"], "label": card["dimension"]} for card in technical_focus_cards],
                    },
                    {
                        "id": "tasks",
                        "label": "任务拆解",
                        "children": [{"id": item["task_name"], "label": item["task_name"]} for item in task_breakdown],
                    },
                    {
                        "id": "meeting",
                        "label": "启动会",
                        "children": [{"id": str(idx), "label": item} for idx, item in enumerate(meeting_agenda, start=1)],
                    },
                    {
                        "id": "ppt",
                        "label": "PPT结构",
                        "children": [{"id": str(item["page"]), "label": item["title"]} for item in ppt_outline],
                    },
                ],
            }
        ],
    }
    return {
        "mode": "mock" if settings.mock_mode else "local_workflow",
        "project_summary": {
            "name": project.name,
            "city": project.city,
            "project_type": project.project_type,
            "phase": project.phase,
            "description": description,
            "knowledge_refs_count": len(refs),
            "summary": f"{project.name} 当前进入项目启动分析，重点是把资料缺口、历史技术复用和会议推进转成可执行任务。",
        },
        "technical_focus_cards": technical_focus_cards,
        "task_breakdown": task_breakdown,
        "meeting_agenda": meeting_agenda,
        "ppt_outline": ppt_outline,
        "risk_list": risk_list,
        "open_questions": open_questions,
        "mindmap_json": mindmap_json,
        "source_refs": refs,
    }


def startup_analysis_to_markdown(payload: dict[str, Any]) -> str:
    summary = payload.get("project_summary", {})
    lines = [
        "# 项目启动分析",
        "",
        f"## 项目摘要",
        f"- 项目：{summary.get('name', '')}",
        f"- 阶段：{summary.get('phase', '')}",
        f"- 结论：{summary.get('summary', '')}",
        "",
        "## 技术重点",
    ]
    for card in payload.get("technical_focus_cards", []):
        lines.append(f"- **{card.get('dimension')}**：{card.get('summary')}")
    lines.extend(["", "## 任务拆解"])
    for task in payload.get("task_breakdown", []):
        lines.append(f"- {task.get('task_name')}（{task.get('owner_role')} / {task.get('priority')}）")
    lines.extend(["", "## 启动会议程"])
    lines.extend(f"{idx}. {item}" for idx, item in enumerate(payload.get("meeting_agenda", []), start=1))
    lines.extend(["", "## PPT结构"])
    for item in payload.get("ppt_outline", []):
        lines.append(f"{item.get('page')}. {item.get('title')}：{item.get('content')}")
    lines.extend(["", "## 风险与未决问题"])
    lines.extend(f"- {item}" for item in payload.get("risk_list", []))
    lines.extend(f"- 待确认：{item}" for item in payload.get("open_questions", []))
    return "\n".join(lines)


def _upsert_startup_skill_card(
    db: Session,
    project: models.Project,
    card_type: str,
    title: str,
    markdown: str,
    output: dict[str, Any],
) -> models.SkillCard:
    card = db.scalar(
        select(models.SkillCard).where(models.SkillCard.project_id == project.id, models.SkillCard.card_type == card_type)
    )
    if not card:
        card = models.SkillCard(project_id=project.id, card_type=card_type)
        db.add(card)
    card.title = title
    card.status = "succeeded"
    card.input_data = json.dumps({"source": "startup_analysis"}, ensure_ascii=False)
    card.output_data = json.dumps(output, ensure_ascii=False)
    card.input_json = json.dumps({"source": "startup_analysis"}, ensure_ascii=False)
    card.output_json = json.dumps(output, ensure_ascii=False)
    card.markdown = markdown
    card.source = "startup_analysis"
    card.created_by = "system"
    card.completed_at = datetime.utcnow()
    return card


def _candidate_markdown(title: str, body: str, sources: list[str]) -> str:
    lines = [
        "---",
        "type: obsidian_candidate",
        f'title: "{title}"',
        "status: pending_review",
        f'updated: "{datetime.utcnow().date().isoformat()}"',
        "---",
        "",
        f"# {title}",
        "",
        body,
        "",
        "## 候选来源",
    ]
    lines.extend(f"- {source}" for source in sources[:8])
    return "\n".join(lines)


def save_startup_analysis(db: Session, project: models.Project, payload: dict[str, Any]) -> models.ProjectReport:
    db.execute(
        delete(models.ProjectReport).where(
            models.ProjectReport.project_id == project.id,
            models.ProjectReport.report_type == "startup_analysis",
        )
    )
    db.execute(delete(models.KnowledgeReference).where(models.KnowledgeReference.project_id == project.id))
    report = models.ProjectReport(
        project_id=project.id,
        report_type="startup_analysis",
        content_json=json.dumps(payload, ensure_ascii=False),
        markdown=startup_analysis_to_markdown(payload),
        model_name=settings.deepseek_model if not settings.mock_mode else "mock",
        mode=payload.get("mode", "mock"),
    )
    db.add(report)
    for ref in payload.get("source_refs", [])[:12]:
        db.add(
            models.KnowledgeReference(
                project_id=project.id,
                source_file=ref.get("source_file", ""),
                source_path=ref.get("source_path", ""),
                chunk_id=ref.get("chunk_id", ""),
                quote=ref.get("quote", ""),
                relevance_score=float(ref.get("relevance_score") or 0),
            )
        )
    existing_task_names = {task.task_name for task in project.tasks}
    for item in payload.get("task_breakdown", []):
        if item["task_name"] in existing_task_names:
            continue
        db.add(
            models.ProjectTask(
                project_id=project.id,
                task_name=item["task_name"],
                task_type=item["task_type"],
                priority=item["priority"],
                owner_role=item["owner_role"],
                estimated_days=int(item["estimated_days"]),
                dependencies=json.dumps(item.get("dependencies", []), ensure_ascii=False),
                risk_level=item["risk_level"],
                status=item["status"],
                output_requirement=item["output_requirement"],
            )
        )
    agenda_markdown = "# 项目启动会议程\n\n" + "\n".join(f"{idx}. {item}" for idx, item in enumerate(payload["meeting_agenda"], start=1))
    if not project.meetings:
        db.add(
            models.Meeting(
                project_id=project.id,
                title=f"{project.name} 项目启动会",
                agenda=agenda_markdown,
                status="scheduled",
                mindmap_json=json.dumps(payload["mindmap_json"], ensure_ascii=False),
                next_actions_json=json.dumps(payload["task_breakdown"], ensure_ascii=False),
            )
        )
    _upsert_startup_skill_card(
        db,
        project,
        "technical_focus",
        "技术重点卡",
        "## 技术重点\n" + "\n".join(f"- **{card['dimension']}**：{card['summary']}" for card in payload["technical_focus_cards"]),
        {"cards": payload["technical_focus_cards"]},
    )
    _upsert_startup_skill_card(
        db,
        project,
        "task_breakdown",
        "任务拆解卡",
        "## 启动任务\n" + "\n".join(f"- {task['task_name']}（{task['owner_role']}）" for task in payload["task_breakdown"]),
        {"items": payload["task_breakdown"]},
    )
    _upsert_startup_skill_card(
        db,
        project,
        "meeting_agenda",
        "会议议程卡",
        agenda_markdown,
        {"items": payload["meeting_agenda"]},
    )
    _upsert_startup_skill_card(
        db,
        project,
        "ppt_outline",
        "PPT结构卡",
        "## PPT结构\n" + "\n".join(f"{item['page']}. {item['title']}：{item['content']}" for item in payload["ppt_outline"]),
        {"slides": payload["ppt_outline"]},
    )
    sources = [ref.get("source_path", "") for ref in payload.get("source_refs", [])]
    db.execute(
        delete(models.KnowledgeItem).where(
            models.KnowledgeItem.project_id == project.id,
            models.KnowledgeItem.item_type == "obsidian_candidate",
        )
    )
    for card in payload.get("technical_focus_cards", []):
        db.add(
            models.KnowledgeItem(
                project_id=project.id,
                source_file=card["title"],
                item_type="obsidian_candidate",
                summary=card["summary"],
                tags=json.dumps(["技术卡", card["dimension"], "候选"], ensure_ascii=False),
                content=_candidate_markdown(card["title"], card["summary"], sources),
            )
        )
    db.commit()
    db.refresh(report)
    return report


def run_startup_analysis(db: Session, project: models.Project) -> dict[str, Any]:
    question = " ".join(
        [
            project.name,
            project.city,
            project.project_type,
            project.phase,
            project.description,
            "日照 退界 面积计算 消防 规划条件 报批 强排 会议纪要",
        ]
    )
    chunks = search_knowledge(db, question, limit=10)
    payload = build_startup_analysis_payload(project, chunks)
    report = save_startup_analysis(db, project, payload)
    return {"report": report, **payload}


def build_team_plan(db: Session, project: models.Project) -> models.TeamPlan:
    employees = list(db.scalars(select(models.DigitalEmployee)))
    tasks = project.tasks
    role_tasks: dict[str, list[str]] = defaultdict(list)
    for task in tasks:
        role_tasks[task.owner_role or "AI项目经理"].append(task.task_name)
    roles = []
    for employee in employees:
        matched = role_tasks.get(employee.name) or role_tasks.get(employee.role) or []
        if matched or employee.name in {"AI项目经理", "AI资料管理员", "AI汇报助理"}:
            roles.append(
                {
                    "name": employee.name,
                    "role": employee.role,
                    "recommended_count": 1,
                    "tasks": matched[:6],
                    "skills": json.loads(employee.skills or "[]"),
                    "intensity": "高" if matched else "中",
                    "risk_note": "需要人工确认任务边界和交付标准。",
                }
            )
    if not roles:
        roles = [
            {
                "name": "AI项目经理",
                "role": "项目负责人 / PM",
                "recommended_count": 1,
                "tasks": [],
                "skills": ["任务拆解"],
                "intensity": "中",
                "risk_note": "Mock 规则生成。",
            }
        ]
    plan = models.TeamPlan(
        project_id=project.id,
        recommended_roles=json.dumps(roles, ensure_ascii=False),
        staffing_summary=f"建议 {len(roles)} 类数字员工参与；需要根据真实任务继续校准投入强度。",
    )
    db.add(plan)
    db.commit()
    db.refresh(plan)
    return plan


def default_meeting_agenda(project: models.Project) -> str:
    return "\n".join(
        [
            f"# {project.name} 项目启动会议程",
            "1. 确认项目背景、业主诉求与当前设计阶段。",
            "2. 核对规划条件、红线、日照、退界、面积计算等技术资料缺口。",
            "3. 复盘历史相似项目可复用经验和风险清单。",
            "4. 明确本轮交付物、节点时间和责任人。",
            "5. 确认下一次业主沟通需要决策的问题。",
        ]
    )


def summarize_meeting(db: Session, meeting: models.Meeting) -> models.Meeting:
    project = db.get(models.Project, meeting.project_id)
    project_name = project.name if project else "当前项目"
    transcript = meeting.transcript.strip()
    basis = transcript or meeting.agenda
    if not basis and project:
        basis = default_meeting_agenda(project)
    actions = [
        {"title": "补齐规划条件、红线、日照和退界资料", "owner": "项目经理", "status": "todo"},
        {"title": "整理历史相似项目技术复用卡", "owner": "AI资料管理员", "status": "todo"},
        {"title": "输出本轮任务拆解和PPT结构", "owner": "AI汇报助理", "status": "todo"},
    ]
    meeting.summary = (
        f"# {project_name} 会议纪要\n\n"
        "## 核心结论\n"
        "- 本轮会议重点围绕项目启动、资料缺口、技术边界和下一步分工展开。\n"
        "- 需要优先确认日照、退界、面积计算、消防/报批等基础约束。\n"
        "- AI 将把会议结论转化为任务看板和后续汇报结构。\n\n"
        "## 原始记录摘要\n"
        f"{basis[:1600] or '暂无会议记录，已按启动会议程生成纪要。'}\n\n"
        "## 下一步\n"
        + "\n".join(f"- {item['title']}（{item['owner']}）" for item in actions)
    )
    meeting.mindmap_json = json.dumps(
        {
            "name": project_name,
            "children": [
                {"name": "资料缺口", "children": [{"name": "规划条件"}, {"name": "日照退界"}, {"name": "面积计算"}]},
                {"name": "任务推进", "children": [{"name": "任务拆解"}, {"name": "责任人"}, {"name": "节点计划"}]},
                {"name": "汇报准备", "children": [{"name": "PPT结构"}, {"name": "风险问题"}, {"name": "业主决策"}]},
            ],
        },
        ensure_ascii=False,
    )
    meeting.next_actions_json = json.dumps(actions, ensure_ascii=False)
    meeting.status = "summarized"
    db.commit()

    existing_names = {task.task_name for task in project.tasks} if project else set()
    if project:
        for item in actions:
            if item["title"] in existing_names:
                continue
            db.add(
                models.ProjectTask(
                    project_id=project.id,
                    task_name=item["title"],
                    task_type="会议待办",
                    priority="高",
                    owner_role=item["owner"],
                    estimated_days=2,
                    dependencies="[]",
                    risk_level="中",
                    status="todo",
                    output_requirement="来自会议纪要自动生成，需要项目经理确认。",
                )
            )
        db.commit()
    db.refresh(meeting)
    return meeting


def run_skill_card(db: Session, project: models.Project, card_type: str, prompt: str = "") -> models.SkillCard:
    card_names = {
        "task_breakdown": "任务拆解卡",
        "technical_focus": "技术重点卡",
        "meeting_minutes": "会议纪要/待办卡",
        "ppt_outline": "PPT结构卡",
    }
    card_type = card_type if card_type in card_names else "task_breakdown"
    tasks = [task.task_name for task in project.tasks[:8]]
    refs = [ref.source_file for ref in project.knowledge_references[:6]]
    meetings = [meeting.title for meeting in project.meetings[:5]]
    outputs: dict[str, Any] = {
        "task_breakdown": {
            "items": tasks
            or ["资料完整性检查", "技术边界复核", "竞品案例收集", "PPT汇报结构搭建"],
            "next": "将任务分配给真实成员或AI数字员工。",
        },
        "technical_focus": {
            "cards": ["日照计算方式", "退界要求", "面积计算方式", "消防/报批风险", "规划条件复核"],
            "knowledge_sources": refs or ["待扫描知识库后自动引用来源"],
        },
        "meeting_minutes": {
            "meetings": meetings or ["项目启动会"],
            "next_actions": ["补齐资料", "确认关键技术边界", "准备下次业主会"],
        },
        "ppt_outline": {
            "slides": ["封面", "项目背景", "技术边界", "历史案例复用", "任务拆解", "风险与下一步"],
            "tone": "面向业主沟通，结论清晰、风险前置、任务可执行。",
        },
    }
    markdown = f"# {card_names[card_type]}\n\n"
    if prompt:
        markdown += f"## 用户需求\n{prompt}\n\n"
    if card_type == "task_breakdown":
        markdown += "## 建议任务\n" + "\n".join(f"- {item}" for item in outputs[card_type]["items"])
    elif card_type == "technical_focus":
        markdown += "## 技术复用重点\n" + "\n".join(f"- {item}" for item in outputs[card_type]["cards"])
        markdown += "\n\n## 来源\n" + "\n".join(f"- {item}" for item in outputs[card_type]["knowledge_sources"])
    elif card_type == "meeting_minutes":
        markdown += "## 会议推进\n" + "\n".join(f"- {item}" for item in outputs[card_type]["next_actions"])
    else:
        markdown += "## PPT框架\n" + "\n".join(f"{idx + 1}. {item}" for idx, item in enumerate(outputs[card_type]["slides"]))

    card = models.SkillCard(
        project_id=project.id,
        card_type=card_type,
        title=card_names[card_type],
        status="succeeded",
        input_json=json.dumps({"prompt": prompt, "project_id": project.id}, ensure_ascii=False),
        output_json=json.dumps(outputs[card_type], ensure_ascii=False),
        markdown=markdown,
        source="mock" if settings.mock_mode else "deepseek",
    )
    db.add(card)
    db.commit()
    db.refresh(card)
    return card
